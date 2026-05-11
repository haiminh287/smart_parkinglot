"""ParkSmart KLTN Slide Generator v8 — Pro design + Speaker notes.

Design principles:
- Title 32pt, body 18pt min, big numbers 60pt+
- 1 idea per slide · breathing room · consistent grid
- Speaker notes for every slide (View → Notes Page in PowerPoint)

Output: docs/slides/ParkSmart_KLTN_Slides.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ─── Color palette ────────────────────────────────────────────────────────
BLUE       = RGBColor(0x1E, 0x40, 0xAF)
BLUE_DEEP  = RGBColor(0x17, 0x2A, 0x6E)
BLUE_SOFT  = RGBColor(0xDB, 0xEA, 0xFE)
BLUE_PALE  = RGBColor(0xEF, 0xF6, 0xFF)
ORANGE     = RGBColor(0xF9, 0x73, 0x16)
ORANGE_SOFT= RGBColor(0xFE, 0xD7, 0xAA)
ORANGE_PALE= RGBColor(0xFF, 0xF7, 0xED)
SLATE      = RGBColor(0x0F, 0x17, 0x2A)
SLATE_2    = RGBColor(0x33, 0x41, 0x55)
GRAY       = RGBColor(0x64, 0x74, 0x8B)
GRAY_SOFT  = RGBColor(0xF1, 0xF5, 0xF9)
GRAY_2     = RGBColor(0xE2, 0xE8, 0xF0)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREEN      = RGBColor(0x10, 0xB9, 0x81)
GREEN_SOFT = RGBColor(0xD1, 0xFA, 0xE5)
PURPLE     = RGBColor(0x8B, 0x5C, 0xF6)
PURPLE_SOFT= RGBColor(0xED, 0xE9, 0xFE)
RED        = RGBColor(0xEF, 0x44, 0x44)
TEAL       = RGBColor(0x14, 0xB8, 0xA6)

FONT       = "Calibri"
FONT_BOLD  = "Calibri"
FONT_MONO  = "Consolas"

SLIDE_W    = Inches(13.333)
SLIDE_H    = Inches(7.5)
TOTAL      = 23

# Typography scale
SIZE_HUGE       = 64   # cover title
SIZE_BIG_NUMBER = 60   # KPI giá trị
SIZE_TITLE      = 32   # slide title
SIZE_SECTION    = 22   # section heading inside slide
SIZE_BODY       = 18   # body text
SIZE_CAPTION    = 14   # caption / supporting
SIZE_LABEL      = 11   # small uppercase labels
SIZE_FOOTER     = 10

# ═══════════════════════════════════════════════════════════════════════════
#  HELPERS
# ═══════════════════════════════════════════════════════════════════════════

def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def set_gradient_bg(slide, hex1, hex2, angle=315):
    """Apply linear gradient bg. PowerPoint requires <p:bg><p:bgPr>...</p:bgPr></p:bg>
    structure inside <p:cSld>, BEFORE <p:spTree>.
    """
    from pptx.oxml.ns import qn
    from lxml import etree
    # Get full <p:sld> root element, find <p:cSld>
    sld = slide._element
    cSld = sld.find(qn('p:cSld'))
    # Remove any existing <p:bg> AND any stray <p:bgPr>
    for tag in ('p:bg', 'p:bgPr'):
        existing = cSld.find(qn(tag))
        if existing is not None:
            cSld.remove(existing)
    # Create <p:bg> as FIRST child of <p:cSld> (must come before <p:spTree>)
    bg = etree.SubElement(cSld, qn('p:bg'))
    spTree = cSld.find(qn('p:spTree'))
    if spTree is not None:
        spTree.addprevious(bg)
    # Build <p:bgPr><a:gradFill>...</a:gradFill></p:bgPr>
    bgPr = etree.SubElement(bg, qn('p:bgPr'))
    grad = etree.SubElement(bgPr, qn('a:gradFill'))
    grad.set('rotWithShape', '1')
    gsLst = etree.SubElement(grad, qn('a:gsLst'))
    for pos, h in [(0, hex1), (100000, hex2)]:
        gs = etree.SubElement(gsLst, qn('a:gs'))
        gs.set('pos', str(pos))
        clr = etree.SubElement(gs, qn('a:srgbClr'))
        clr.set('val', h.lstrip('#').upper())
    lin = etree.SubElement(grad, qn('a:lin'))
    lin.set('ang', str(int(angle * 60000)))
    lin.set('scaled', '1')
    # Required tail: bgRef or empty for proper schema
    # Actually <p:bg> only needs <p:bgPr> OR <p:bgRef>, we have bgPr ✓




def add_rect(slide, left, top, width, height, fill, line=None,
             shape=MSO_SHAPE.RECTANGLE):
    box = slide.shapes.add_shape(shape, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    if line is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line
        box.line.width = Pt(0.75)
    box.shadow.inherit = False
    return box


def add_text(slide, left, top, width, height, text,
             size=SIZE_BODY, bold=False, color=SLATE, align=PP_ALIGN.LEFT,
             font=FONT, anchor=MSO_ANCHOR.TOP, italic=False, line_spacing=1.2):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(50000)
    tf.margin_top = tf.margin_bottom = Emu(20000)
    tf.vertical_anchor = anchor
    lines = [text] if isinstance(text, str) else list(text)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_section_label(slide, num, label, accent=BLUE):
    """Top-left section number + label."""
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(8), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = f"{num:02d}"
    r1.font.name = FONT_BOLD
    r1.font.size = Pt(SIZE_LABEL)
    r1.font.bold = True
    r1.font.color.rgb = accent
    r2 = p.add_run()
    r2.text = f"   {label.upper()}"
    r2.font.name = FONT
    r2.font.size = Pt(SIZE_LABEL)
    r2.font.color.rgb = GRAY


LOGO_PATH = str(Path(__file__).parent / "assets" / "hcmou_logo.png")


def add_logo(slide, mode="small"):
    """Add HCMOU university logo.
    mode='small' → top-right corner (content slides, 0.75")
    mode='large' → cover-style (1.0")
    """
    if not Path(LOGO_PATH).exists():
        return
    if mode == "small":
        # Top-right corner — bigger for KLTN report standard
        slide.shapes.add_picture(LOGO_PATH,
            Inches(12.1), Inches(0.2),
            height=Inches(0.9))
    elif mode == "large":
        slide.shapes.add_picture(LOGO_PATH,
            Inches(0.5), Inches(0.35),
            height=Inches(1.5))


def add_footer(slide, slide_num):
    add_text(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
             "ParkSmart  ·  Khóa luận tốt nghiệp 2026  ·  Nguyễn Hải Minh",
             size=SIZE_FOOTER, color=GRAY)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3),
             f"{slide_num:02d} / {TOTAL:02d}",
             size=SIZE_FOOTER, color=GRAY, bold=True,
             align=PP_ALIGN.RIGHT, font=FONT_BOLD)


def add_title(slide, text, color=SLATE, top=Inches(0.95), accent=ORANGE):
    add_text(slide, Inches(0.5), top, Inches(12.3), Inches(0.75),
             text, size=SIZE_TITLE, bold=True, color=color, font=FONT_BOLD,
             line_spacing=1.0)
    # Accent underline
    add_rect(slide, Inches(0.5), top + Inches(0.78), Inches(0.8), Inches(0.06),
             accent)


def add_notes(slide, text):
    """Add speaker notes — visible in PowerPoint Notes Page view."""
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def add_card(slide, left, top, width, height, accent=BLUE):
    """White card with left accent stripe + subtle border."""
    add_rect(slide, left, top, width, height, WHITE, line=GRAY_2)
    add_rect(slide, left, top, Inches(0.12), height, accent)


def add_kpi_block(slide, left, top, width, height, value, label,
                   accent=ORANGE, value_size=44, label_size=12):
    """Big number top, label bottom — non-overlapping zones."""
    add_rect(slide, left, top, width, height, GRAY_SOFT)
    add_rect(slide, left, top, width, Inches(0.06), accent)
    val_h = height * 0.62
    add_text(slide, left, top + Inches(0.1), width, val_h,
             value, size=value_size, bold=True, color=accent,
             align=PP_ALIGN.CENTER, font=FONT_BOLD, line_spacing=1.0,
             anchor=MSO_ANCHOR.MIDDLE)
    lbl_top = top + Inches(0.1) + val_h
    lbl_h = height - Inches(0.15) - val_h
    add_text(slide, left, lbl_top, width, lbl_h,
             label, size=label_size, color=SLATE,
             align=PP_ALIGN.CENTER, line_spacing=1.1,
             anchor=MSO_ANCHOR.MIDDLE)


def add_image_safe(slide, image_path, left, top, width=None, height=None):
    p = Path(image_path)
    if not p.exists():
        # placeholder rectangle if image missing
        w = width or Inches(8)
        h = height or Inches(4)
        add_rect(slide, left, top, w, h, GRAY_SOFT)
        add_text(slide, left, top + h / 2 - Inches(0.2), w, Inches(0.4),
                 f"[ Image not found: {p.name} ]", size=12, italic=True,
                 color=GRAY, align=PP_ALIGN.CENTER)
        return None
    if width and height:
        return slide.shapes.add_picture(str(p), left, top, width=width, height=height)
    if width:
        return slide.shapes.add_picture(str(p), left, top, width=width)
    if height:
        return slide.shapes.add_picture(str(p), left, top, height=height)
    return slide.shapes.add_picture(str(p), left, top)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════════════

def slide_01(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # Diagonal gradient: dark navy → medium navy (professional, không trùng logo)
    set_gradient_bg(s, "#0A1330", "#1E3A5F", angle=315)
    # Diagonal accent block top-right
    add_rect(s, Inches(11.2), Inches(0), Inches(2.13), Inches(1.5), ORANGE)
    # Bottom strip orange
    add_rect(s, Inches(0), Inches(7.3), SLIDE_W, Inches(0.2), ORANGE)

    # Logo HCMOU — BIGGER (1.5") on WHITE rounded card top-left
    add_rect(s, Inches(0.4), Inches(0.25), Inches(1.7), Inches(1.7),
             WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_logo(s, mode="large")  # logo 1.5" — see helper

    # University name beside logo (positioned for bigger logo)
    add_text(s, Inches(2.3), Inches(0.55), Inches(8.5), Inches(0.4),
             "TRƯỜNG ĐẠI HỌC MỞ THÀNH PHỐ HỒ CHÍ MINH",
             size=16, bold=True, color=WHITE)
    add_text(s, Inches(2.3), Inches(0.98), Inches(8.5), Inches(0.35),
             "KHOA CÔNG NGHỆ THÔNG TIN",
             size=12, color=ORANGE_SOFT)
    add_text(s, Inches(2.3), Inches(1.33), Inches(8.5), Inches(0.35),
             "NGÀNH KHOA HỌC MÁY TÍNH",
             size=12, color=ORANGE_SOFT)

    # Top-right: thesis label (on top of orange diagonal)
    add_text(s, Inches(11.4), Inches(0.5), Inches(1.9), Inches(0.4),
             "KLTN · 2026",
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Center section — title block
    add_rect(s, Inches(0), Inches(2.55), SLIDE_W, Inches(0.06), ORANGE)
    add_text(s, Inches(0.7), Inches(2.7), Inches(12), Inches(0.35),
             "ĐỀ TÀI KHÓA LUẬN",
             size=12, bold=True, color=ORANGE)

    # Big project name
    add_text(s, Inches(0.7), Inches(3.1), Inches(12), Inches(1.4),
             "PARKSMART", size=80, bold=True,
             color=WHITE, font=FONT_BOLD, line_spacing=1.0)

    # Sub-title (chỉ "tích hợp IoT" — không có "AI và Digital Twin")
    add_text(s, Inches(0.7), Inches(4.55), Inches(12), Inches(0.6),
             "Hệ thống quản lý Bãi đỗ xe Thông minh",
             size=26, bold=True, color=ORANGE_SOFT, font=FONT_BOLD)
    add_text(s, Inches(0.7), Inches(5.1), Inches(12), Inches(0.6),
             "tích hợp IoT",
             size=26, bold=True, color=ORANGE_SOFT, font=FONT_BOLD)

    # Divider
    add_rect(s, Inches(0.7), Inches(5.95), Inches(2.5), Inches(0.05), ORANGE)

    # Author block — clean & professional (positioned below subtitle)
    add_rect(s, Inches(0.45), Inches(6.25), Inches(8.5), Inches(0.9),
             SLATE_2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(0.7), Inches(6.32), Inches(8), Inches(0.28),
             "SINH VIÊN THỰC HIỆN",
             size=10, bold=True, color=ORANGE)
    add_text(s, Inches(0.7), Inches(6.6), Inches(8), Inches(0.48),
             "Nguyễn Hải Minh",
             size=22, bold=True, color=WHITE, font=FONT_BOLD)
    add_text(s, Inches(5.2), Inches(6.6), Inches(3.7), Inches(0.48),
             "MSSV 2251012093",
             size=18, color=ORANGE_SOFT, italic=True,
             anchor=MSO_ANCHOR.MIDDLE)

    add_notes(s,
        "Xin chào quý hội đồng, em là Nguyễn Hải Minh, MSSV 2251012093, "
        "Khoa Công nghệ Thông tin, ngành Khoa học Máy tính. Khóa luận hôm nay "
        "em trình bày đề tài: Hệ thống quản lý Bãi đỗ xe Thông minh tích hợp IoT, AI và Digital Twin. "
        "Em sẽ trình bày trong khoảng 15-20 phút, gồm các phần: đặt vấn đề, "
        "kiến trúc hệ thống, 4 mô-đun AI, IoT phần cứng, Unity Digital Twin, "
        "kết quả thực nghiệm, phân tích thực tế và hướng phát triển. Em xin bắt đầu."
    )


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — ĐẶT VẤN ĐỀ
# ═══════════════════════════════════════════════════════════════════════════

def slide_02(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(s, "#FFFFFF", "#EFF6FF", angle=315)
    add_section_label(s, 1, "Đặt vấn đề")
    add_logo(s)
    add_title(s, "Hạn chế của bãi đỗ xe truyền thống")

    # 5 horizontal cards
    pains = [
        ("01", "Ùn tắc cửa", "30-60s", "mỗi xe ra/vào", RED),
        ("02", "Sai sót tiền", "≈5%", "giao dịch nhầm mệnh giá", ORANGE),
        ("03", "Nhân sự", "21-35tr", "/tháng cho 3-5 bảo vệ", PURPLE),
        ("04", "Không đặt trước", "60%", "user phải đi 2-3 bãi", BLUE),
        ("05", "Không analytics", "Sổ tay", "không tối ưu được", SLATE_2),
    ]
    card_w = Inches(2.4)
    card_h = Inches(3.5)
    gap = Inches(0.15)
    total_w = card_w * 5 + gap * 4
    start_left = (SLIDE_W - total_w) / 2
    top = Inches(2.4)
    for i, (num, title, big, sub, c) in enumerate(pains):
        left = start_left + (card_w + gap) * i
        # Card
        add_rect(s, left, top, card_w, card_h, GRAY_SOFT)
        # Top color band with number
        add_rect(s, left, top, card_w, Inches(0.6), c)
        add_text(s, left, top + Inches(0.12), card_w, Inches(0.45),
                 num, size=24, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_BOLD)
        # Title
        add_text(s, left + Inches(0.15), top + Inches(0.85),
                 card_w - Inches(0.3), Inches(0.5),
                 title, size=18, bold=True, color=SLATE,
                 align=PP_ALIGN.CENTER, font=FONT_BOLD)
        # BIG stat
        add_text(s, left + Inches(0.15), top + Inches(1.5),
                 card_w - Inches(0.3), Inches(1.1),
                 big, size=42, bold=True, color=c,
                 align=PP_ALIGN.CENTER, font=FONT_BOLD,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        # Sub
        add_text(s, left + Inches(0.2), top + Inches(2.7),
                 card_w - Inches(0.4), card_h - Inches(2.8),
                 sub, size=13, color=SLATE_2,
                 align=PP_ALIGN.CENTER, line_spacing=1.3)

    # Bottom takeaway
    add_rect(s, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.6), BLUE)
    add_text(s, Inches(0.5), Inches(6.18), Inches(12.3), Inches(0.55),
             "→  Cần giải pháp tự động hóa toàn diện bằng AI",
             size=18, bold=True, color=WHITE, italic=True,
             align=PP_ALIGN.CENTER, font=FONT_BOLD,
             anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, 2)
    add_notes(s,
        "Em đã khảo sát 3 bãi đỗ xe nội đô: Vincom Q1, The Sun Avenue và Saigon Centre. "
        "Em phát hiện 5 nỗi đau chung. Thứ nhất, ùn tắc cửa — mỗi xe mất 30-60 giây "
        "khi check-in thủ công, lên tới 60-90 giây khi thanh toán cash, gây kẹt xe "
        "200m vào giờ cao điểm. Thứ hai, khoảng 5% giao dịch tiền mặt bị nhầm mệnh giá "
        "— bãi 1000 vé/tháng thất thoát 3-7 triệu. Thứ ba, nhân sự thủ công 3-5 bảo vệ "
        "tốn 21-35 triệu/tháng cố định, không scale theo doanh thu. Thứ tư, 60% user "
        "phải đi 2-3 bãi mới tìm được chỗ. Thứ năm, không có analytics, quản lý bằng "
        "sổ tay, không tối ưu được pricing. Vì vậy, cần giải pháp tự động hóa toàn diện."
    )


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — MỤC TIÊU & PHẠM VI
# ═══════════════════════════════════════════════════════════════════════════

def slide_03(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(s, "#FFFFFF", "#EFF6FF", angle=315)
    add_section_label(s, 2, "Mục tiêu & Phạm vi")
    add_logo(s)
    add_title(s, "Mục tiêu hệ thống — 4 chức năng cốt lõi")

    # 4 outcome-focused goal cards (no tech names — đó là slide 4)
    goals = [
        ("Tự động kiểm soát ra/vào",
         "Camera đọc biển số → tự động mở barrier\nGiảm 30-60s wait time mỗi xe · không cần bảo vệ",
         BLUE),
        ("Tự động thanh toán tiền mặt",
         "Quét tờ tiền → nhận diện 9 mệnh giá VN\nGiảm 5% sai sót tính tiền · tự động cộng phí",
         ORANGE),
        ("Trợ lý ảo hỗ trợ 24/7",
         "Chatbot trả lời FAQ + guide đặt chỗ\nGiảm tải support center · trả lời từ docs thật",
         PURPLE),
        ("Mô phỏng + test toàn hệ thống",
         "Digital Twin 3D mô phỏng bãi xe\nTest E2E không cần phần cứng · tiết kiệm 50tr setup",
         GREEN),
    ]
    cw = Inches(5.8)
    ch = Inches(1.85)
    gap = Inches(0.18)
    g_left = Inches(0.5)
    g_top = Inches(2.4)
    for i, (t, body, c) in enumerate(goals):
        col = i % 2
        row = i // 2
        left = g_left + (cw + gap) * col
        top = g_top + (ch + gap) * row
        add_card(s, left, top, cw, ch, accent=c)
        add_text(s, left + Inches(0.35), top + Inches(0.2),
                 cw - Inches(0.5), Inches(0.55),
                 t, size=20, bold=True, color=c, font=FONT_BOLD)
        add_text(s, left + Inches(0.35), top + Inches(0.85),
                 cw - Inches(0.5), ch - Inches(1.0),
                 body, size=15, color=SLATE_2, line_spacing=1.4)

    # (Removed "production-ready" tagline per user feedback)
    add_footer(s, 3)
    add_notes(s,
        "Hệ thống ParkSmart hướng đến 4 chức năng cốt lõi giải quyết 5 nỗi đau ở slide trước. "
        "Một, tự động kiểm soát ra vào: camera đọc biển số xe, hệ thống tự động mở barrier "
        "khi khớp booking — giảm thời gian wait từ 30-60 giây xuống còn vài giây, không cần "
        "bảo vệ check thủ công. "
        "Hai, tự động thanh toán tiền mặt: máy quét tờ tiền, AI nhận diện 9 mệnh giá VN, "
        "tự động cộng vào running total — loại bỏ 5% sai sót do nhân viên đếm nhầm. "
        "Ba, trợ lý ảo 24/7: chatbot trả lời FAQ + guide user đặt chỗ — giảm tải call center, "
        "trả lời từ knowledge base thật có citation source. "
        "Bốn, mô phỏng và test toàn hệ thống: Digital Twin 3D cho phép test E2E mà không "
        "cần phần cứng thật — tiết kiệm 50 triệu chi phí prototype hardware. "
        "Ở slide tiếp theo em sẽ giải thích HOW — chọn công nghệ gì để hiện thực hoá 4 chức năng này."
    )


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — TECH STACK
# ═══════════════════════════════════════════════════════════════════════════

def slide_04(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(s, "#FFFFFF", "#EFF6FF", angle=315)
    add_section_label(s, 3, "Công nghệ sử dụng")
    add_logo(s)
    add_title(s, "6 quyết định công nghệ quan trọng nhất")

    # 6 decisions with detailed WHY (not just listing)
    decisions = [
        ("01", "Microservices",
         "Tách AI service (GPU) khỏi service nhẹ. Fault isolation: AI lỗi không sụp booking. Scale từng service độc lập theo tải.",
         BLUE),
        ("02", "Python + Go kết hợp",
         "Python cho AI (PyTorch ecosystem mạnh). Go cho Gateway (latency thấp, concurrent connections cao). Mỗi việc 1 công cụ.",
         PURPLE),
        ("03", "Gemini 2.5 Flash (LLM)",
         "Tiếng Việt support tốt · 200× rẻ hơn GPT-4 ($0.15 vs $30/1M tokens) · free tier 1500 req/day đủ cho KLTN.",
         ORANGE),
        ("04", "EfficientNetV2-S (Vision)",
         "Pretrained ImageNet 84% · Fused-MBConv train 4× nhanh hơn V1 · fit GTX 1650 4GB với mixed precision fp16.",
         GREEN),
        ("05", "ChromaDB (Vector store cho RAG)",
         "Embedded vector DB — không cần external service như Pinecone. Local persistent, deploy đơn giản trong Docker Compose.",
         TEAL),
        ("06", "Unity Digital Twin",
         "C# scripting dễ học · NavMesh AI sẵn cho xe tự lái · validate IoT contract end-to-end mà không cần phần cứng thật.",
         RED),
    ]
    cw = Inches(4.0)
    ch = Inches(2.05)
    gap_x = Inches(0.15)
    gap_y = Inches(0.2)
    grid_w = cw * 3 + gap_x * 2
    start_left = (SLIDE_W - grid_w) / 2
    start_top = Inches(2.4)
    for i, (num, name, why, c) in enumerate(decisions):
        col = i % 3
        row = i // 3
        left = start_left + (cw + gap_x) * col
        top = start_top + (ch + gap_y) * row
        add_card(s, left, top, cw, ch, accent=c)
        # Number badge top-right
        add_rect(s, left + cw - Inches(0.55), top + Inches(0.12),
                 Inches(0.4), Inches(0.4), c, shape=MSO_SHAPE.OVAL)
        add_text(s, left + cw - Inches(0.55), top + Inches(0.12),
                 Inches(0.4), Inches(0.4), num,
                 size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_BOLD,
                 anchor=MSO_ANCHOR.MIDDLE)
        # Tech name
        add_text(s, left + Inches(0.25), top + Inches(0.18),
                 cw - Inches(0.85), Inches(0.45),
                 name, size=16, bold=True, color=c, font=FONT_BOLD)
        # Why explanation (detail)
        add_text(s, left + Inches(0.25), top + Inches(0.7),
                 cw - Inches(0.45), ch - Inches(0.8),
                 why, size=12, color=SLATE_2, line_spacing=1.4)

    add_footer(s, 4)
    add_notes(s,
        "Em sẽ giải thích 6 quyết định công nghệ quan trọng nhất, mỗi cái có lý do cụ thể. "
        "Một, Microservices: tách AI service cần GPU mạnh khỏi service nhẹ, fault isolation. "
        "Hai, kết hợp Python với Go: Python cho AI vì ecosystem mạnh, Go cho Gateway vì latency thấp + concurrent tốt. "
        "Ba, chọn Gemini 2.5 Flash không phải GPT-4: tiếng Việt support tốt, rẻ hơn 200 lần, free tier đủ KLTN. "
        "Bốn, chọn EfficientNetV2-S không phải ResNet: pretrained ImageNet 84%, Fused-MBConv train 4 lần nhanh, fit GPU 4GB. "
        "Năm, ChromaDB cho RAG: embedded vector DB không cần external service như Pinecone, deploy đơn giản. "
        "Sáu, Unity Digital Twin: C# dễ học, NavMesh AI sẵn cho xe tự lái, validate IoT contract mà không cần phần cứng."
    )


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — KIẾN TRÚC TỔNG THỂ (image)
# ═══════════════════════════════════════════════════════════════════════════

def slide_05(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(s, "#FFFFFF", "#EFF6FF", angle=315)
    add_section_label(s, 4, "Kiến trúc tổng thể")
    add_logo(s)
    add_title(s, "Kiến trúc hệ thống ParkSmart")

    # ── LEFT: Architecture image (width-constrained) ────────────────
    img_left = Inches(0.5)
    img_top = Inches(2.0)
    img_w = Inches(7.4)
    img_h_max = Inches(4.4)
    pic = add_image_safe(s, r"C:/Users/MINH/Desktop/overview_project.jpg",
                         img_left, img_top, width=img_w)
    if pic and pic.height and pic.height > img_h_max:
        aspect = pic.width / pic.height
        pic.height = img_h_max
        pic.width = Emu(int(img_h_max * aspect))

    # Caption under image
    cap_y = (pic.top + pic.height) if pic else (img_top + img_h_max)
    add_text(s, img_left, cap_y + Inches(0.05),
             img_w, Inches(0.3),
             "Hình 1. Sơ đồ kiến trúc tổng thể hệ thống ParkSmart",
             size=11, italic=True, color=GRAY,
             align=PP_ALIGN.CENTER)

    # ── RIGHT: Explanation panel — 5 layers ─────────────────────────
    panel_left = Inches(8.1)
    panel_top = Inches(2.0)
    panel_w = Inches(4.75)
    add_text(s, panel_left, panel_top, panel_w, Inches(0.35),
             "GIẢI THÍCH SƠ ĐỒ KIẾN TRÚC",
             size=11, bold=True, color=ORANGE, font=FONT_BOLD)
    explanations = [
        ("L1", "Client", "Web app · Unity simulator · ESP32 IoT", BLUE),
        ("L2", "Gateway (Go)", "Xác thực session · phát khóa bảo mật · rate limit", ORANGE),
        ("L3", "Microservices", "10 service: auth · booking · parking · AI · chatbot", PURPLE),
        ("L4", "Event Bus", "RabbitMQ fan-out · 1 event → nhiều consumer async", GREEN),
        ("L5", "Datastore", "MySQL relational + Redis × 6 DBs cô lập", TEAL),
    ]
    e_top = panel_top + Inches(0.45)
    e_h = Inches(0.78)
    e_gap = Inches(0.04)
    for i, (badge, name, desc, c) in enumerate(explanations):
        top = e_top + (e_h + e_gap) * i
        add_rect(s, panel_left, top, panel_w, e_h, WHITE, line=GRAY_2)
        add_rect(s, panel_left, top, Inches(0.55), e_h, c)
        add_text(s, panel_left, top, Inches(0.55), e_h,
                 badge, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_BOLD,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, panel_left + Inches(0.65), top + Inches(0.08),
                 panel_w - Inches(0.7), Inches(0.3),
                 name, size=13, bold=True, color=c, font=FONT_BOLD)
        add_text(s, panel_left + Inches(0.65), top + Inches(0.38),
                 panel_w - Inches(0.7), Inches(0.38),
                 desc, size=10, color=SLATE_2, line_spacing=1.25)

    # (Bottom callout removed — explanation panel on right covers it)
    add_footer(s, 5)
    add_notes(s,
        "Đây là sơ đồ kiến trúc tổng thể của hệ thống ParkSmart. Em chia làm 7 layer. "
        "Layer 1 là Client gồm Web SPA, Unity Digital Twin và ESP32 IoT. Layer 2 là "
        "Cloudflare Tunnel cho HTTPS public. Layer 3 là API Gateway viết bằng Go, "
        "xác thực session và rate limit. Layer 4 là 10 microservices: auth, booking, "
        "parking, vehicle, notification, payment, ai, chatbot. Layer 5 là RabbitMQ "
        "làm event bus async fan-out. Layer 6 là Realtime Go cho WebSocket push. "
        "Layer 7 là datastore: MySQL 8 và Redis chia làm 6 DBs. "
        "Điểm chốt là chỉ có 1 entry point public — mọi service nội bộ chỉ chấp nhận "
        "request đã được Gateway xác thực qua khóa bảo mật."
    )


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — SO SÁNH KIẾN TRÚC
# ═══════════════════════════════════════════════════════════════════════════

def slide_06(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(s, "#FFFFFF", "#EFF6FF", angle=315)
    add_section_label(s, 5, "Lựa chọn kiến trúc")
    add_logo(s)
    add_title(s, "Vì sao chọn Microservices?")

    # 3-column comparison
    cols = [
        ("MONOLITH", "1 codebase\n1 deploy unit\nDev nhanh ban đầu",
         ["✓ Simple", "✗ Khó scale", "✗ Sụp toàn bộ", "✗ Khó đa stack"], GRAY),
        ("MICROSERVICES", "10 services độc lập\nGiao tiếp HTTP + Event\nGateway pattern",
         ["✓ Scale từng phần", "✓ Fault isolation", "✓ Đa stack dễ", "⚠ Setup phức tạp"], BLUE),
        ("SERVERLESS", "Function-as-a-service\nAuto-scaling\nPay per request",
         ["✓ Auto scale", "⚠ Cold start", "⚠ Cost cao", "⚠ Vendor lock-in"], PURPLE),
    ]
    cw = Inches(3.85)
    ch = Inches(4.0)
    gap = Inches(0.3)
    grid_w = cw * 3 + gap * 2
    start_left = (SLIDE_W - grid_w) / 2
    top = Inches(2.4)
    for i, (name, desc, pros, c) in enumerate(cols):
        left = start_left + (cw + gap) * i
        is_chosen = (i == 1)
        # Highlight chosen
        add_rect(s, left, top, cw, ch, BLUE_PALE if is_chosen else WHITE,
                 line=BLUE if is_chosen else GRAY_2)
        # Header band
        add_rect(s, left, top, cw, Inches(0.6), c)
        add_text(s, left, top + Inches(0.12), cw, Inches(0.45),
                 name, size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_BOLD)
        # Description
        add_text(s, left + Inches(0.25), top + Inches(0.85),
                 cw - Inches(0.4), Inches(1.2),
                 desc, size=14, color=SLATE_2, line_spacing=1.4,
                 align=PP_ALIGN.CENTER)
        # Pros list — bounded height to leave room for badge
        add_text(s, left + Inches(0.35), top + Inches(2.05),
                 cw - Inches(0.5), Inches(1.3),
                 pros, size=14, color=SLATE, line_spacing=1.5)
        # Chosen badge
        if is_chosen:
            add_rect(s, left + Inches(0.4), top + ch - Inches(0.55),
                     cw - Inches(0.8), Inches(0.4), BLUE)
            add_text(s, left, top + ch - Inches(0.55), cw, Inches(0.4),
                     "ĐÃ CHỌN", size=14, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, font=FONT_BOLD,
                     anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, 6)
    add_notes(s,
        "Em đã so sánh 3 kiến trúc. Monolith dễ làm ban đầu nhưng khó scale, fault một "
        "service làm sụp cả hệ thống, khó dùng đa stack. Serverless auto-scale tốt nhưng "
        "có cold start, cost cao và vendor lock-in. Em chọn Microservices vì 4 lý do: "
        "Một, scale được từng service riêng — AI service cần GPU mạnh, các service khác "
        "thì nhẹ, tách ra rồi mới scale từng phần. Hai, fault isolation — nếu AI service "
        "lỗi thì booking và payment vẫn hoạt động. Ba, dễ dùng đa stack — Python cho AI, "
        "Go cho gateway latency thấp, JS cho frontend. Bốn, thể hiện được nhiều skill "
        "kiến trúc cho khóa luận. Trade-off là setup ban đầu phức tạp hơn, nhưng đáng giá."
    )


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — EVENT-DRIVEN + REDIS
# ═══════════════════════════════════════════════════════════════════════════

def slide_07(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(s, "#FFFFFF", "#EFF6FF", angle=315)
    add_section_label(s, 6, "Patterns kiến trúc")
    add_logo(s)
    add_title(s, "Event-Driven + Cache layer")

    # Left: RabbitMQ fan-out
    add_text(s, Inches(0.5), Inches(2.3), Inches(7.5), Inches(0.4),
             "RABBITMQ — FAN-OUT EVENT BUS",
             size=14, bold=True, color=PURPLE, font=FONT_BOLD)

    # Source box
    add_rect(s, Inches(0.5), Inches(2.85), Inches(2.8), Inches(0.7), BLUE)
    add_text(s, Inches(0.5), Inches(3.0), Inches(2.8), Inches(0.5),
             "booking-service", size=16, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font=FONT_BOLD,
             anchor=MSO_ANCHOR.MIDDLE)
    # Arrow
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                               Inches(3.4), Inches(3.05), Inches(0.4), Inches(0.3))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ORANGE
    arrow.line.fill.background()
    # Hub
    add_rect(s, Inches(3.95), Inches(2.85), Inches(2.0), Inches(0.7), PURPLE)
    add_text(s, Inches(3.95), Inches(3.0), Inches(2.0), Inches(0.5),
             "RabbitMQ", size=16, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font=FONT_BOLD,
             anchor=MSO_ANCHOR.MIDDLE)

    # 4 consumers
    consumers = [("notification", "email + push"),
                 ("payment", "QR code"),
                 ("analytics", "log event"),
                 ("chatbot", "user state")]
    cw_c = Inches(1.7)
    cgap = Inches(0.07)
    cstart = Inches(0.55)
    cy = Inches(4.4)
    for i, (name, sub) in enumerate(consumers):
        cleft = cstart + (cw_c + cgap) * i
        # Down arrow
        ar = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                cleft + cw_c / 2 - Inches(0.1),
                                Inches(3.6), Inches(0.2), Inches(0.6))
        ar.fill.solid(); ar.fill.fore_color.rgb = PURPLE
        ar.line.fill.background()
        # Consumer
        add_rect(s, cleft, cy, cw_c, Inches(0.8), GRAY_SOFT, line=PURPLE)
        add_text(s, cleft, cy + Inches(0.08), cw_c, Inches(0.35),
                 name, size=13, bold=True, color=PURPLE,
                 align=PP_ALIGN.CENTER, font=FONT_BOLD)
        add_text(s, cleft, cy + Inches(0.43), cw_c, Inches(0.32),
                 sub, size=11, color=GRAY,
                 align=PP_ALIGN.CENTER, italic=True)

    # Right: Redis 6 DBs
    add_text(s, Inches(8.5), Inches(2.3), Inches(4.3), Inches(0.4),
             "REDIS — 6 DBs CÔ LẬP",
             size=14, bold=True, color=GREEN, font=FONT_BOLD)
    redis_items = [
        ("DB 0", "Celery — task queue"),
        ("DB 1", "auth + gateway session"),
        ("DB 2", "booking cache"),
        ("DB 3", "parking slot status"),
        ("DB 4", "vehicle service"),
        ("DB 5", "realtime WS pub/sub"),
        ("DB 6", "chatbot state TTL 30p"),
    ]
    r_top = Inches(2.85)
    r_h = Inches(0.55)
    for i, (k, v) in enumerate(redis_items):
        t = r_top + r_h * i
        bg = WHITE if i % 2 == 0 else GRAY_SOFT
        add_rect(s, Inches(8.5), t, Inches(4.3), r_h, bg)
        add_rect(s, Inches(8.5), t, Inches(0.1), r_h, GREEN)
        add_text(s, Inches(8.7), t + Inches(0.12),
                 Inches(0.8), Inches(0.35),
                 k, size=14, bold=True, color=GREEN, font=FONT_BOLD,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(9.55), t + Inches(0.12),
                 Inches(3.2), Inches(0.35),
                 v, size=13, color=SLATE,
                 anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, 7)
    add_notes(s,
        "Em dùng 2 pattern quan trọng. Một là RabbitMQ làm event bus theo fan-out pattern. "
        "Ví dụ: booking-service phát event booking.created một lần, RabbitMQ broadcast tới "
        "4 consumer cùng lúc — notification gửi email, payment generate QR, analytics log "
        "event, chatbot cập nhật state. Lợi ích là async fire-and-forget, loose coupling, "
        "thêm consumer mới không cần sửa publisher. Hai là Redis tách thành 6 DBs riêng "
        "cho từng service: DB 0 cho Celery task queue, DB 1 cho auth + gateway session, "
        "DB 2 cho booking cache, DB 3 cho parking slot, DB 5 cho realtime WebSocket pub/sub, "
        "DB 6 cho chatbot state với TTL 30 phút. Cô lập key-space giúp flush DB lẻ "
        "không ảnh hưởng service khác."
    )


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — END-TO-END FLOW
# ═══════════════════════════════════════════════════════════════════════════

def slide_08(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(s, "#FFFFFF", "#EFF6FF", angle=315)
    add_section_label(s, 7, "Luồng hoạt động")
    add_logo(s)
    add_title(s, "Một chu trình hoàn chỉnh — 3 phase")

    phases = [
        ("PHASE 1", "ĐẶT CHỖ ONLINE",
         "User mở web/app  →  Gateway xác thực  →  Booking-service tạo booking  →  MySQL persist  →  RabbitMQ event  →  Email + QR confirmation",
         BLUE),
        ("PHASE 2", "CHECK-IN VÀO BÃI",
         "Xe đến cổng  →  Camera ANPR chụp biển  →  AI đọc text  →  Match booking (Levenshtein ≤ 1)  →  ESP32 mở barrier  →  Dashboard cập nhật",
         GREEN),
        ("PHASE 3", "CHECK-OUT + THANH TOÁN",
         "Xe ra  →  Cash scanner ESP32  →  AI nhận diện 9 mệnh giá  →  Tính phí  →  Payment đóng booking  →  Barrier mở  →  Email biên lai",
         ORANGE),
    ]
    p_top = Inches(2.3)
    p_h = Inches(1.4)
    p_gap = Inches(0.15)
    for i, (label, title, body, c) in enumerate(phases):
        top = p_top + (p_h + p_gap) * i
        # Card
        add_rect(s, Inches(0.5), top, Inches(12.3), p_h, GRAY_SOFT)
        # Left phase block
        add_rect(s, Inches(0.5), top, Inches(2.5), p_h, c)
        add_text(s, Inches(0.5), top + Inches(0.25), Inches(2.5), Inches(0.4),
                 label, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_BOLD)
        add_text(s, Inches(0.5), top + Inches(0.65), Inches(2.5), Inches(0.6),
                 title, size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_BOLD,
                 line_spacing=1.1)
        # Right body
        add_text(s, Inches(3.2), top + Inches(0.2), Inches(9.4), p_h - Inches(0.4),
                 body, size=15, color=SLATE, line_spacing=1.6,
                 anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, 8)
    add_notes(s,
        "Đây là luồng hoạt động end-to-end của hệ thống — em chia làm 3 phase. "
        "Phase 1, Đặt chỗ online: User mở web hoặc app, Gateway xác thực session, "
        "Booking-service tạo booking, persist xuống MySQL, phát event qua RabbitMQ, "
        "Notification gửi email kèm QR confirmation. "
        "Phase 2, Check-in vào bãi: Xe đến cổng, camera ANPR chụp biển số, AI đọc text, "
        "match với booking đã đăng ký bằng Levenshtein distance tối đa 1, ESP32 mở barrier, "
        "dashboard admin cập nhật realtime. "
        "Phase 3, Check-out + Thanh toán: Xe đến cổng ra, cash scanner ESP32 chụp tờ tiền, "
        "AI nhận diện 9 mệnh giá, cộng tổng đến khi đủ phí, Payment-service đóng booking, "
        "barrier mở, gửi email biên lai. Đây là full flow."
    )


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — AI #1: PLATE
# ═══════════════════════════════════════════════════════════════════════════

def _ai_module_slide(prs, slide_num, section_num, title_text, accent,
                     pipeline, kpis, techniques, notes_text):
    """Generic AI module slide layout."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(s, "#FFFFFF", "#EFF6FF", angle=315)
    add_section_label(s, section_num, "AI Module", accent=accent)
    add_logo(s)
    add_title(s, title_text)

    # Pipeline (top, full-width) — compact
    add_text(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(0.3),
             "PIPELINE", size=12, bold=True, color=accent, font=FONT_BOLD)
    pipe_top = Inches(2.2)
    pipe_h = Inches(0.65)
    n = len(pipeline)
    arrow_w = Inches(0.22)
    box_w = (Inches(12.3) - arrow_w * (n - 1) - Inches(0.05) * (n - 1)) / n
    for i, step in enumerate(pipeline):
        left = Inches(0.5) + (box_w + arrow_w + Inches(0.05)) * i
        add_rect(s, left, pipe_top, box_w, pipe_h, GRAY_SOFT)
        add_rect(s, left, pipe_top, box_w, Inches(0.05), accent)
        add_text(s, left + Inches(0.08), pipe_top + Inches(0.1),
                 box_w - Inches(0.16), pipe_h - Inches(0.15),
                 step, size=12, bold=True, color=SLATE,
                 align=PP_ALIGN.CENTER, font=FONT_BOLD,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
        if i < n - 1:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    left + box_w + Inches(0.025),
                                    pipe_top + Inches(0.28), arrow_w, Inches(0.2))
            ar.fill.solid(); ar.fill.fore_color.rgb = accent
            ar.line.fill.background()

    # 4 KPI blocks — moved up
    kpi_top = Inches(3.05)
    kpi_h = Inches(0.95)
    kw = (Inches(12.3) - Inches(0.15) * 3) / 4
    for i, (val, lbl) in enumerate(kpis):
        left = Inches(0.5) + (kw + Inches(0.15)) * i
        add_kpi_block(s, left, kpi_top, kw, kpi_h, val, lbl, accent=accent)

    # Techniques — 3 cols × 2 rows · BIG cards (1.3" tall — room for full text)
    add_text(s, Inches(0.5), Inches(4.15), Inches(12.3), Inches(0.3),
             "KỸ THUẬT ÁP DỤNG  ·  Vì sao chọn các kỹ thuật này?",
             size=12, bold=True, color=accent, font=FONT_BOLD)
    t_top = Inches(4.5)
    card_h = Inches(1.2)
    card_w = (Inches(12.3) - Inches(0.2) * 2) / 3
    for i, (name, desc) in enumerate(techniques[:6]):
        col = i % 3
        row = i // 3
        left = Inches(0.5) + (card_w + Inches(0.2)) * col
        top = t_top + (card_h + Inches(0.1)) * row
        # Card body
        add_rect(s, left, top, card_w, card_h, WHITE, line=GRAY_2)
        # Top accent stripe
        add_rect(s, left, top, card_w, Inches(0.06), accent)
        # Tech name
        add_text(s, left + Inches(0.2), top + Inches(0.12),
                 card_w - Inches(0.3), Inches(0.32),
                 name, size=13, bold=True, color=accent, font=FONT_BOLD)
        # Description body — bigger area, 11pt font
        add_text(s, left + Inches(0.2), top + Inches(0.46),
                 card_w - Inches(0.3), card_h - Inches(0.52),
                 desc, size=11, color=SLATE_2, line_spacing=1.3)

    add_footer(s, slide_num)
    add_notes(s, notes_text)


def slide_09(prs):
    _ai_module_slide(
        prs, slide_num=9, section_num=8,
        title_text="AI #1 — Nhận diện biển số xe",
        accent=BLUE,
        pipeline=["YOLOv8\nDetect", "Crop +\nWarp", "TrOCR\nPrimary",
                  "EasyOCR\nFallback", "Levenshtein\nMatch"],
        kpis=[("96.3%", "Detection accuracy"),
              ("92.1%", "OCR accuracy"),
              ("94.5%", "End-to-end"),
              ("<0.5%", "False accept rate")],
        techniques=[
            ("YOLOv8 Detection",
             "Anchor-free + decoupled head. Train 2000 ảnh biển VN đạt 96%. Stable hơn v11 mới ra, ecosystem matured."),
            ("Perspective Warp",
             "Nắn 4 góc bbox biển nghiêng thành chữ nhật thẳng. Deterministic không cần ML. +5% OCR accuracy."),
            ("TrOCR Primary",
             "Vision Transformer Microsoft pretrained. Mạnh nhất với biển khó (mờ, nghiêng, đêm) — 96% hard cases."),
            ("EasyOCR + Tesseract",
             "Fallback khi TrOCR conf < 0.6. Triple-engine cover edge cases. Tổng accuracy 92-96%."),
            ("Levenshtein ≤ 1 Match",
             "Tolerate sai 1 ký tự (OCR hay nhầm O↔0, I↔1). UX tốt mà vẫn an toàn — false accept < 0.5%."),
            ("Pre-warm Lifespan",
             "Load AI models vào RAM khi FastAPI startup. Tránh cold-start 20s. Trade-off: tốn ~2GB RAM."),
        ],
        notes_text=(
            "Module AI thứ nhất là nhận diện biển số xe. Em dùng pipeline 2-stage. "
            "Stage 1 là YOLOv8 detect bbox biển — em chọn YOLOv8 vì anchor-free và decoupled head, "
            "đạt 96% accuracy so với classical CV chỉ 60%. Sau đó perspective warp nắn biển "
            "nghiêng thành chữ nhật thẳng, tăng OCR accuracy thêm 5%. Stage 2 là OCR cascade 3-engine: "
            "TrOCR Vision Transformer của Microsoft đọc đầu tiên, mạnh nhất với biển khó như mờ, "
            "nghiêng, đêm — đạt 96% trên hard cases. Nếu confidence thấp, fallback EasyOCR và "
            "Tesseract. Cuối cùng so OCR với booking bằng Levenshtein distance, cho phép sai tối đa "
            "1 ký tự để tolerate nhầm O với 0, I với 1. Kết quả end-to-end là 94.5% accuracy, "
            "false accept rate dưới 0.5%, latency khoảng 2 giây. Em cũng pre-warm tất cả model "
            "ngay khi service khởi động để tránh cold-start 20 giây."
        )
    )


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — AI #2: BANKNOTE
# ═══════════════════════════════════════════════════════════════════════════

def slide_10(prs):
    _ai_module_slide(
        prs, slide_num=10, section_num=9,
        title_text="AI #2 — Nhận diện mệnh giá tiền VN  ·  UNIQUE",
        accent=ORANGE,
        pipeline=["T1 Quality\nLaplacian", "T2 Detector\nHSV+Contour",
                  "T3 Color\nGabor+SVM", "T4 Deep\nENetV2-S+TTA"],
        kpis=[("98.22%", "Accuracy 9 mệnh giá"),
              ("≥99.5%", "Precision-at-accept"),
              ("85-90%", "Accept rate"),
              ("~50ms", "Latency trung bình")],
        techniques=[
            ("Cascade 4-tầng",
             "Quality → Detector → Color SVM → Deep. 70% dừng ở tầng 3, TB 50ms. Đánh đổi latency vs accuracy."),
            ("EfficientNetV2-S",
             "Pretrained ImageNet 84%. Fused-MBConv train 4× nhanh hơn V1. Fit GTX 1650 4GB với mixed precision."),
            ("TTA × 5 biến thể",
             "Predict 5 augmented versions → average softmax. +1-2% accuracy không cần train thêm."),
            ("Rejection Margin-based",
             "Top1≥0.92 AND margin≥0.25 → accept. Precision-first cho giao dịch tiền — 'thà từ chối còn hơn sai'."),
            ("Siamese + OneClass SVM",
             "Bank-grade 3-stage. Siamese verify cặp embedding. OneClass SVM detect anomaly tiền giả/lỗi."),
            ("Weighted Sampling",
             "Class imbalance 200k:457 vs 1800. Sampler weight 1/count + Albumentations augment realistic."),
        ],
        notes_text=(
            "Đây là điểm độc đáo của em — AI nhận diện 9 mệnh giá tiền Việt Nam, "
            "chưa bãi xe nào tại Việt Nam có. Em dùng pipeline cascade 4 tầng từ nhẹ đến nặng. "
            "Tầng 1 quality check bằng Laplacian variance để filter ảnh mờ, mất 2ms. "
            "Tầng 2 detector dùng HSV mask và contour, mất 5ms. "
            "Tầng 3 color classifier Gabor filter cộng SVM xử lý 70% case dễ trong 5ms. "
            "Nếu chưa chắc chắn, tầng 4 mới dùng deep model EfficientNetV2-S cộng TTA 5 biến thể, "
            "mất 150ms. Trung bình tổng chỉ 50ms, nhanh hơn 3 lần so với chạy deep model cho mọi case. "
            "Em còn áp dụng rejection margin-based: top-1 confidence cao chưa đủ, phải kèm margin "
            "với top-2 — không đủ tự tin thì từ chối, yêu cầu user scan lại. Đây là triết lý "
            "precision-first: thà từ chối còn hơn sai cho giao dịch tiền mặt. Bank-grade pipeline "
            "còn có Siamese network và OneClass SVM cho verification và phát hiện tiền giả. "
            "Kết quả accuracy 98.22%, precision-at-accept ≥ 99.5%."
        )
    )


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — AI #3: CHATBOT
# ═══════════════════════════════════════════════════════════════════════════

def slide_11(prs):
    """Use AI helper for consistent layout with other AI slides."""
    _ai_module_slide(
        prs, slide_num=11, section_num=10,
        title_text="AI #3 — Trợ lý ảo tiếng Việt  ·  RAG-powered",
        accent=PURPLE,
        pipeline=["Preproc\nnormalize", "Intent\n16 class",
                  "Entity\nJSON schema", "Action\ncall API", "Response\nformat"],
        kpis=[("16", "Intents > 93% acc"),
              ("0.8s", "Latency/message"),
              ("RAG", "Citation từ docs"),
              ("Free", "Gemini tier 1500/day")],
        techniques=[
            ("DDD 3-Layer",
             "Domain · Application · Infrastructure tách rõ. Đổi LLM chỉ sửa Infrastructure, không động Domain."),
            ("Gemini 2.5 Flash",
             "Tiếng Việt tốt · 200× rẻ hơn GPT-4 ($0.15 vs $30/1M tokens) · latency 0.8s · free tier."),
            ("Schema-driven Entity",
             "Mỗi intent có JSON schema riêng. LLM bắt buộc trả structured. Tránh free-text parsing."),
            ("Hybrid Confidence",
             "0.4×LLM + 0.3×Entity + 0.3×Context. LLM tự tin một mình không đủ — gộp 3 dim chính xác hơn."),
            ("Booking Wizard Multi-turn",
             "State machine 3 bước (vehicle → lot → slot). Redis state TTL 30 phút. 1 câu không đủ."),
            ("RAG + ChromaDB",
             "ChromaDB vector store + multilingual embeddings. Top-K=3 cosine ≥ 0.35. FAQ có citation."),
        ],
        notes_text=(
            "Module AI thứ ba: chatbot trợ lý ảo tiếng Việt với kiến trúc DDD 3-layer. "
            "Domain layer chứa intent enum 16 class, Application orchestration logic, "
            "Infrastructure swap được LLM provider. Em chọn Gemini 2.5 Flash vì hỗ trợ tiếng Việt tốt, "
            "rẻ hơn GPT-4 200 lần, free tier 1500 req/day. Schema-driven entity extraction đảm bảo "
            "structured output. Hybrid Confidence kết hợp 3 dimension chống false confident. "
            "Booking Wizard guide user step-by-step qua 3 bước. Đặc biệt RAG với ChromaDB vector store "
            "+ sentence-transformers multilingual — chatbot trả lời FAQ từ knowledge base thật, có "
            "citation source, giảm hallucination. Kết quả: 16 intents accuracy > 93%, latency 0.8s."
        )
    )



# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — AI #4: SLOT DETECTION
# ═══════════════════════════════════════════════════════════════════════════

def slide_12(prs):
    _ai_module_slide(
        prs, slide_num=12, section_num=11,
        title_text="AI #4 — Nhận diện ô đỗ xe",
        accent=TEAL,
        pipeline=["HSV Mask\nviền cam", "Morphology\nCLOSE+DILATE",
                  "findContours\nRETR_CCOMP", "Classify\nHSV patch"],
        kpis=[("100%", "Bãi trống detect"),
              ("86-93%", "Bãi đầy detect"),
              ("<80ms", "Latency/frame"),
              ("5 FPS", "Realtime dashboard")],
        techniques=[
            ("Classical CV (không YOLO)",
             "Thử YOLO fail vì Unity primitive ≠ COCO. HSV + contour: deterministic, 100% bãi trống."),
            ("Morphology CLOSE+DILATE",
             "Anti-alias làm viền slot hở. CLOSE kernel 3×3 + DILATE 1 iter giúp contour khép kín."),
            ("findContours RETR_CCOMP",
             "Tìm 'lỗ' bên trong grid cam (mỗi ô = 1 lỗ). CCOMP cho 2-level hierarchy outer+holes — đúng nhất."),
            ("Status HSV Patch",
             "Sample HSV trung tâm 35% ô → phân loại màu. Xe có thể che viền nhưng tâm thường visible."),
            ("Row Reconstruction",
             "Xe to che viền giữa 2 ô → mất 2 ô. Cluster theo Y, interpolate slot bị che."),
            ("Best-Grid Cache",
             "Memo grid lớn nhất từng detect. Frame kém → fallback cache. Stability 86-93% bãi đầy."),
        ],
        notes_text=(
            "Module AI thứ tư là nhận diện ô đỗ xe. Em đã thử YOLO fine-tune nhưng fail vì "
            "Unity primitive khác với COCO dataset. Em chuyển sang Classical CV, kết quả "
            "100% accuracy khi bãi trống, 86-93% khi bãi đầy. Pipeline: HSV mask viền cam "
            "của slot, morphology CLOSE và DILATE để khép kín viền bị anti-alias, "
            "findContours với flag RETR_CCOMP để tìm 2-level hierarchy — viền ngoài và lỗ "
            "bên trong, lỗ chính là ô đỗ. Cuối cùng classify status bằng cách sample HSV patch "
            "trung tâm 35% của ô. Em còn có 2 kỹ thuật quan trọng: row reconstruction để "
            "interpolate slot bị xe che, và best-grid cache để memo grid lớn nhất từng detect "
            "được — khi frame hiện tại kém, fallback về cache. Latency dưới 80ms mỗi frame, "
            "realtime 5 FPS trên dashboard admin."
        )
    )


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — AI PROOF (confusion matrix image)
# ═══════════════════════════════════════════════════════════════════════════

def slide_13(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(s, "#FFFFFF", "#EFF6FF", angle=315)
    add_section_label(s, 12, "Kết quả thực nghiệm", accent=GREEN)
    add_logo(s)
    add_title(s, "AI hoạt động trên dữ liệu thật")

    # Left: 4 AI summary
    items = [
        ("Plate OCR", "Test 500 ảnh biển VN", "94.5%", "End-to-end accuracy", BLUE),
        ("Banknote", "Test 1500 ảnh holdout", "98.22%", "9 mệnh giá VN", ORANGE),
        ("Slot Detection", "Test Unity scenes", "86-100%", "Theo occupancy", TEAL),
        ("Chatbot", "16 intents × 50 câu", ">93%", "Intent accuracy", PURPLE),
    ]
    cw = Inches(3.4)
    ch = Inches(0.95)
    cy = Inches(2.4)
    for i, (name, dataset, val, lbl, c) in enumerate(items):
        top = cy + (ch + Inches(0.1)) * i
        add_rect(s, Inches(0.5), top, cw + Inches(2), ch, WHITE, line=GRAY_2)
        add_rect(s, Inches(0.5), top, Inches(0.12), ch, c)
        add_text(s, Inches(0.7), top + Inches(0.1),
                 Inches(2.5), Inches(0.4),
                 name, size=16, bold=True, color=c, font=FONT_BOLD)
        add_text(s, Inches(0.7), top + Inches(0.5),
                 Inches(2.5), Inches(0.5),
                 dataset, size=11, color=GRAY, italic=True)
        # Big metric
        add_text(s, Inches(3.5), top + Inches(0.05),
                 Inches(2.0), Inches(0.6),
                 val, size=32, bold=True, color=c,
                 align=PP_ALIGN.LEFT, font=FONT_BOLD,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(3.5), top + Inches(0.65),
                 Inches(2.0), Inches(0.35),
                 lbl, size=11, color=SLATE_2, italic=True)

    # Right: confusion matrix
    add_text(s, Inches(7.5), Inches(2.4), Inches(5.3), Inches(0.35),
             "CONFUSION MATRIX BANKNOTE", size=12, bold=True, color=ORANGE,
             font=FONT_BOLD, align=PP_ALIGN.CENTER)
    add_image_safe(s,
        r"C:/Users/MINH/Documents/Zalo_Received_Files/Project_Main/docs/screenshots/09-banknote-confusion-matrix.png",
        Inches(7.5), Inches(2.8), width=Inches(5.3))

    # Bottom methodology — fits above footer
    add_rect(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4), GREEN_SOFT)
    add_text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4),
             "Train/Val/Test 70/15/15  ·  Holdout test set  ·  Cross-validation  ·  So sánh với baseline",
             size=12, bold=True, color=GREEN,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, 13)
    add_notes(s,
        "Đây là kết quả thực nghiệm của 4 AI module em đo trên test set độc lập. "
        "AI Plate OCR test 500 ảnh biển Việt Nam tự thu thập, end-to-end accuracy 94.5%. "
        "AI Banknote test 1500 ảnh holdout — 15% dataset không seen during training, "
        "accuracy 9 mệnh giá đạt 98.22%. Slot Detection test trên Unity scenes ngẫu nhiên, "
        "accuracy 86-100% tùy theo occupancy. Chatbot test 16 intents với khoảng 50 câu "
        "mỗi intent, intent accuracy trên 93%. Bên phải là confusion matrix của banknote — "
        "có thể thấy class 200k và 500k hay nhầm nhau nhất do màu đỏ-tím gần nhau. "
        "Rejection logic của em bù đắp bằng cách yêu cầu user scan lại khi không đủ confidence. "
        "Phương pháp em dùng là split train/val/test 70-15-15, holdout test set không seen "
        "during training, cross-validation cho banknote, và so sánh với baseline MobileNetV3 "
        "đạt 89% và Classical CV đạt 60%."
    )


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — IoT HARDWARE
# ═══════════════════════════════════════════════════════════════════════════

def slide_14(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(s, "#FFFFFF", "#EFF6FF", angle=315)
    add_section_label(s, 13, "IoT Hardware", accent=TEAL)
    add_logo(s)
    add_title(s, "4 ESP32 Edge Devices  +  BOM ~26.5 triệu (chi phí thực tế)")

    # Left: 4 ESP32 cards (2x2)
    devices = [
        ("GATE-IN-01", "Cổng vào", "Camera 4K + Servo + LED", BLUE),
        ("GATE-OUT-01", "Cổng ra", "Camera 4K + LED status", GREEN),
        ("VERIFY-SLOT-01", "Verify ô đỗ", "Camera 2MP overview", PURPLE),
        ("CASH-PAY-01", "Thanh toán cash", "Cash scanner + LED + Buzzer", ORANGE),
    ]
    cw = Inches(3.85)
    ch = Inches(1.85)
    g_left = Inches(0.5)
    g_top = Inches(2.4)
    gap = Inches(0.15)
    for i, (devid, role, hw, c) in enumerate(devices):
        col = i % 2
        row = i // 2
        left = g_left + (cw + gap) * col
        top = g_top + (ch + gap) * row
        add_rect(s, left, top, cw, ch, WHITE, line=GRAY_2)
        add_rect(s, left, top, cw, Inches(0.5), c)
        add_text(s, left + Inches(0.2), top + Inches(0.07),
                 cw - Inches(0.3), Inches(0.4),
                 devid, size=14, bold=True, color=WHITE, font=FONT_MONO,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, left + Inches(0.25), top + Inches(0.65),
                 cw - Inches(0.4), Inches(0.4),
                 role, size=18, bold=True, color=SLATE, font=FONT_BOLD)
        add_text(s, left + Inches(0.25), top + Inches(1.1),
                 cw - Inches(0.4), Inches(0.7),
                 hw, size=13, color=SLATE_2, line_spacing=1.4)

    # Right: BOM
    bom_left = Inches(8.7)
    add_text(s, bom_left, Inches(2.4), Inches(4.3), Inches(0.35),
             "BOM PHẦN CỨNG", size=12, bold=True, color=ORANGE, font=FONT_BOLD)
    bom = [
        ("Cameras (2× IP 4K + 3× dome 2MP)", "9.0tr"),
        ("ESP32 × 4 + Servo + sensors", "1.5tr"),
        ("LED status (per device)", "0.5tr"),
        ("NVR 4ch + Switch PoE 8-port", "4.5tr"),
        ("Mini-PC N100/i5 (used) + Coral USB", "8.0tr"),
        ("Phụ kiện (tủ, dây, khung)", "3.0tr"),
    ]
    b_top = Inches(2.85)
    b_h = Inches(0.42)
    for i, (k, v) in enumerate(bom):
        t = b_top + (b_h + Inches(0.04)) * i
        bg = WHITE if i % 2 == 0 else GRAY_SOFT
        add_rect(s, bom_left, t, Inches(4.3), b_h, bg, line=GRAY_2)
        add_text(s, bom_left + Inches(0.15), t + Inches(0.06),
                 Inches(2.9), b_h - Inches(0.1),
                 k, size=12, color=SLATE,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, bom_left + Inches(3.05), t + Inches(0.06),
                 Inches(1.1), b_h - Inches(0.1),
                 v, size=13, bold=True, color=ORANGE,
                 align=PP_ALIGN.RIGHT, font=FONT_BOLD,
                 anchor=MSO_ANCHOR.MIDDLE)
    # Total
    total_top = b_top + (b_h + Inches(0.04)) * 6 + Inches(0.05)
    add_rect(s, bom_left, total_top, Inches(4.3), Inches(0.55), BLUE_DEEP)
    add_text(s, bom_left + Inches(0.15), total_top, Inches(2), Inches(0.55),
             "TỔNG", size=14, bold=True, color=ORANGE,
             font=FONT_BOLD, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, bom_left + Inches(2.2), total_top, Inches(2.0), Inches(0.55),
             "~26.5tr", size=22, bold=True, color=WHITE,
             align=PP_ALIGN.RIGHT, font=FONT_BOLD,
             anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, 14)
    add_notes(s,
        "Phần IoT em thiết kế 4 ESP32 edge devices. "
        "GATE-IN-01 ở cổng vào: camera IP 4K, loop sensor phát hiện xe, button manual, "
        "servo MG996R cho barrier, LED status. "
        "GATE-OUT-01 ở cổng ra: camera 4K, display LCD hiển thị phí, servo và LED. "
        "VERIFY-SLOT-01 cho verify ô đỗ: camera dome 2MP overview tầng V1, "
        "LED strip WS2812 cho 50 slots. "
        "CASH-PAY-01 cho thanh toán cash: camera scanner focus tờ tiền, LED và buzzer. "
        "Bên phải là BOM chi tiết — cameras 9 triệu, ESP32 và phụ kiện 1.5 triệu, "
        "LED status 0.5 triệu, NVR và network 4.5 triệu, Mini-PC N100 cộng Coral USB 8 triệu, "
        "phụ kiện 3 triệu. Tổng cộng khoảng 26.5 triệu cho bãi 50 chỗ. "
        "Em chọn ESP32 vì có Wi-Fi tích hợp, ADC, servo control, giá chỉ 250k. "
        "Raspberry Pi đắt gấp 5 lần, Arduino không có Wi-Fi sẵn."
    )


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 15 — IoT COMMUNICATION
# ═══════════════════════════════════════════════════════════════════════════

def slide_15(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(s, "#FFFFFF", "#EFF6FF", angle=315)
    add_section_label(s, 14, "IoT giao tiếp", accent=TEAL)
    add_logo(s)
    add_title(s, "Cách thức ESP32 giao tiếp với hệ thống")

    flows = [
        ("01", "KHI XE ĐẾN CỔNG VÀO",
         "Cảm biến phát hiện xe → camera ESP32 chụp ảnh biển → AI đọc text → đối chiếu booking → mở barrier",
         BLUE),
        ("02", "KHI XE ĐỖ VÀO Ô",
         "Camera tầng V1 chụp định kỳ → AI nhận diện slot có xe → ESP32 đổi LED xanh sang đỏ → cập nhật dashboard",
         GREEN),
        ("03", "KHI XE RA + THANH TOÁN",
         "Cash scanner ESP32 chụp tờ tiền → AI phân loại 9 mệnh giá → cộng tổng phí → mở barrier ra → email biên lai",
         ORANGE),
    ]
    f_top = Inches(2.4)
    f_h = Inches(1.1)
    for i, (num, title, body, c) in enumerate(flows):
        top = f_top + (f_h + Inches(0.15)) * i
        add_rect(s, Inches(0.5), top, Inches(12.3), f_h, GRAY_SOFT)
        add_rect(s, Inches(0.5), top, Inches(0.8), f_h, c)
        add_text(s, Inches(0.5), top + Inches(0.3), Inches(0.8), Inches(0.5),
                 num, size=28, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_BOLD,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.55), top + Inches(0.12),
                 Inches(11), Inches(0.4),
                 title, size=16, bold=True, color=c, font=FONT_BOLD)
        add_text(s, Inches(1.55), top + Inches(0.55),
                 Inches(11), f_h - Inches(0.65),
                 body, size=14, color=SLATE, line_spacing=1.4)

    # Bottom 2 callouts
    cb_top = Inches(6.0)
    cb_h = Inches(0.85)
    cw = Inches(6.0)
    add_rect(s, Inches(0.5), cb_top, cw, cb_h, PURPLE_SOFT)
    add_rect(s, Inches(0.5), cb_top, Inches(0.1), cb_h, PURPLE)
    add_text(s, Inches(0.7), cb_top + Inches(0.08), cw - Inches(0.3), Inches(0.3),
             "XÁC THỰC THIẾT BỊ", size=11, bold=True, color=PURPLE, font=FONT_BOLD)
    add_text(s, Inches(0.7), cb_top + Inches(0.35), cw - Inches(0.3), cb_h - Inches(0.4),
             "Mỗi ESP32 có khóa định danh duy nhất. Backend nhận diện đúng thiết bị · giả mạo bị từ chối.",
             size=13, color=SLATE, line_spacing=1.3)

    add_rect(s, Inches(6.8), cb_top, cw, cb_h, GREEN_SOFT)
    add_rect(s, Inches(6.8), cb_top, Inches(0.1), cb_h, GREEN)
    add_text(s, Inches(7.0), cb_top + Inches(0.08), cw - Inches(0.3), Inches(0.3),
             "GIÁM SÁT SỨC KHỎE (HEARTBEAT)", size=11, bold=True, color=GREEN, font=FONT_BOLD)
    add_text(s, Inches(7.0), cb_top + Inches(0.35), cw - Inches(0.3), cb_h - Inches(0.4),
             "Báo cáo trạng thái mỗi 30s. Mất tín hiệu > 90s → cảnh báo offline · bảo vệ manual override.",
             size=13, color=SLATE, line_spacing=1.3)

    add_footer(s, 15)
    add_notes(s,
        "ESP32 giao tiếp với hệ thống qua 3 luồng tương tác chính. "
        "Luồng 1, khi xe đến cổng vào: Cảm biến phát hiện xe, camera ESP32 chụp ảnh biển số, "
        "gửi về AI service đọc text, đối chiếu với booking đã đăng ký. Nếu khớp thì mở barrier. "
        "Luồng 2, khi xe đỗ vào ô: Camera tổng quan tầng V1 chụp định kỳ, AI nhận diện slot có xe, "
        "ESP32 đổi LED từ xanh sang đỏ và cập nhật realtime lên admin dashboard. "
        "Luồng 3, khi xe ra và thanh toán: Cash scanner ESP32 chụp tờ tiền, AI phân loại 9 mệnh giá, "
        "cộng tổng phí, khi đủ thì mở barrier ra và gửi email biên lai. "
        "Về xác thực: mỗi ESP32 được cấp khóa định danh duy nhất khi cài đặt — backend nhận diện "
        "đúng thiết bị, giả mạo bị từ chối ngay. Về giám sát sức khỏe: ESP32 báo cáo trạng thái "
        "mỗi 30 giây, mất tín hiệu > 90s thì cảnh báo offline và bảo vệ có thể manual override."
    )


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 16 — UNITY DIGITAL TWIN
# ═══════════════════════════════════════════════════════════════════════════

def slide_16(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(s, "#FFFFFF", "#EFF6FF", angle=315)
    add_section_label(s, 15, "Digital Twin", accent=GREEN)
    add_logo(s)
    add_title(s, "Unity Digital Twin  ·  Triển khai và Mở rộng đa bãi")

    # Top: 3 deployment modes
    add_text(s, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.35),
             "3 CHẾ ĐỘ TRIỂN KHAI", size=12, bold=True, color=GREEN, font=FONT_BOLD)
    modes = [
        ("DEV", "Editor Play Mode", "Hot-reload C# · debug realtime · dev iteration", BLUE),
        ("EXE", "Standalone Build", ".exe Windows / .app Mac · không cần Unity install", ORANGE),
        ("CI",  "Headless Mode", "-batchmode -nographics · GitHub Actions E2E nightly", PURPLE),
    ]
    cw = Inches(4.0)
    ch = Inches(1.65)
    gap = Inches(0.15)
    g_left = (SLIDE_W - cw * 3 - gap * 2) / 2
    g_top = Inches(2.7)
    for i, (badge, title, desc, c) in enumerate(modes):
        left = g_left + (cw + gap) * i
        add_rect(s, left, g_top, cw, ch, WHITE, line=GRAY_2)
        add_rect(s, left, g_top, Inches(0.85), ch, c)
        add_text(s, left, g_top + Inches(0.5), Inches(0.85), Inches(0.5),
                 badge, size=20, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_BOLD,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, left + Inches(1.0), g_top + Inches(0.2),
                 cw - Inches(1.15), Inches(0.4),
                 title, size=16, bold=True, color=SLATE, font=FONT_BOLD)
        add_text(s, left + Inches(1.0), g_top + Inches(0.65),
                 cw - Inches(1.15), ch - Inches(0.75),
                 desc, size=12, color=SLATE_2, line_spacing=1.4)

    # Bottom: 4 lot models supported
    add_text(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.35),
             "MULTI-LOT FLEXIBILITY  ·  4 MÔ HÌNH BÃI ĐÃ HỖ TRỢ",
             size=12, bold=True, color=GREEN, font=FONT_BOLD)
    models = [
        ("M1", "Linear single-floor", "30-50 chỗ", "Bãi nhỏ, gần cửa", BLUE),
        ("M2", "Block grid + zones", "100-200 chỗ", "TTTM lớn", PURPLE),
        ("M3", "Multi-floor stacked", "3-5 tầng", "Tòa văn phòng", ORANGE),
        ("M4", "Mixed (xe+EV+motor)", "Mọi loại xe", "Bãi tổng hợp", GREEN),
    ]
    mw = Inches(3.0)
    mh = Inches(1.65)
    mg = Inches(0.07)
    m_left = (SLIDE_W - mw * 4 - mg * 3) / 2
    m_top = Inches(5.1)
    for i, (badge, title, scale, use, c) in enumerate(models):
        left = m_left + (mw + mg) * i
        add_rect(s, left, m_top, mw, mh, WHITE, line=GRAY_2)
        add_rect(s, left, m_top, mw, Inches(0.4), c)
        add_text(s, left, m_top + Inches(0.06), mw, Inches(0.3),
                 badge, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_BOLD)
        add_text(s, left + Inches(0.15), m_top + Inches(0.5),
                 mw - Inches(0.3), Inches(0.4),
                 title, size=14, bold=True, color=SLATE, font=FONT_BOLD)
        add_text(s, left + Inches(0.15), m_top + Inches(0.95),
                 mw - Inches(0.3), Inches(0.3),
                 scale, size=12, color=c, italic=True, font=FONT_BOLD)
        add_text(s, left + Inches(0.15), m_top + Inches(1.25),
                 mw - Inches(0.3), Inches(0.35),
                 use, size=11, color=GRAY)

    add_footer(s, 16)
    add_notes(s,
        "Unity Digital Twin của em là điểm độc đáo — chưa đồ án KLTN VN nào có. "
        "Em hỗ trợ 3 chế độ triển khai. "
        "Mode 1, Editor Play Mode: chạy trong Unity Editor với hot-reload C# scripts, "
        "debug logs realtime, dùng cho dev iteration. "
        "Mode 2, Standalone Build: build ra file .exe Windows hoặc .app macOS, "
        "không cần cài Unity, deliver demo qua USB hoặc zip. "
        "Mode 3, Headless CI Mode: chạy với flag -batchmode -nographics, "
        "GitHub Actions chạy E2E test nightly, regression test simulator vs backend, "
        "không cần GPU hay display. "
        "Về flexibility, em thiết kế config-driven: backend trả về JSON layout với floors, "
        "zones, slot grid; ParkingLotGenerator.cs đọc config và tạo procedural scene runtime. "
        "Hiện em đã hỗ trợ 4 mô hình bãi: M1 Linear single-floor cho bãi nhỏ 30-50 chỗ; "
        "M2 Block grid với zones cho TTTM 100-200 chỗ; M3 Multi-floor stacked cho tòa "
        "văn phòng 3-5 tầng; M4 Mixed cho bãi tổng hợp có cả xe hơi, xe máy và xe điện."
    )




# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 18 — KHÓ KHĂN
# ═══════════════════════════════════════════════════════════════════════════

def slide_18(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(s, "#FFFFFF", "#EFF6FF", angle=315)
    add_section_label(s, 16, "Khó khăn & Bài học")
    add_logo(s)
    add_title(s, "Failures cũng là kinh nghiệm")

    issues = [
        ("YOLO fail trên slot", "Unity primitive ≠ COCO", "Fallback OpenCV HSV", RED),
        ("Banknote accuracy thấp", "Backbone yếu, không reject", "EfficientNetV2-S + TTA + Reject", ORANGE),
        ("GTX 1650 4GB OOM", "fp32 + batch lớn", "Mixed precision fp16 + batch 12", PURPLE),
        ("EasyOCR cold-start 20s", "Lazy load model", "Pre-warm trong lifespan()", BLUE),
        ("Gemini quota limit", "Free tier 1500/day", "Dual backend SDK + proxy", GREEN),
        ("Class imbalance 200k", "Sampler thiên majority", "WeightedSampler + class_weight", TEAL),
    ]
    cw = Inches(4.0)
    ch = Inches(1.45)
    gap = Inches(0.15)
    grid_w = cw * 3 + gap * 2
    g_left = (SLIDE_W - grid_w) / 2
    g_top = Inches(2.4)
    for i, (problem, cause, fix, c) in enumerate(issues):
        col = i % 3
        row = i // 3
        left = g_left + (cw + gap) * col
        top = g_top + (ch + gap) * row
        add_rect(s, left, top, cw, ch, WHITE, line=GRAY_2)
        # Top red strip
        add_rect(s, left, top, cw, Inches(0.4), c)
        add_text(s, left + Inches(0.2), top + Inches(0.05),
                 cw - Inches(0.3), Inches(0.3),
                 problem, size=13, bold=True, color=WHITE, font=FONT_BOLD)
        # Cause
        add_text(s, left + Inches(0.2), top + Inches(0.5),
                 cw - Inches(0.3), Inches(0.4),
                 "Nguyên nhân: " + cause, size=11, color=GRAY, italic=True)
        # Fix
        add_text(s, left + Inches(0.2), top + Inches(0.95),
                 cw - Inches(0.3), Inches(0.45),
                 "→ " + fix, size=12, bold=True, color=GREEN, font=FONT_BOLD,
                 line_spacing=1.3)

    # Bottom takeaway
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
             "7 sprints  ·  ~200 commits  ·  refactor 4 god-classes về < 300 dòng",
             size=14, italic=True, color=GRAY,
             align=PP_ALIGN.CENTER)
    add_footer(s, 17)
    add_notes(s,
        "Suốt quá trình em gặp 6 khó khăn lớn, mỗi cái em rút ra bài học. "
        "Một, YOLO fail trên slot detection vì Unity primitive khác COCO dataset — "
        "em fallback sang OpenCV HSV và RETR_CCOMP, đạt 100% bãi trống. "
        "Hai, Banknote accuracy v1 chỉ 89% — em redesign với EfficientNetV2-S, "
        "TTA và rejection logic, lên 98.22%. "
        "Ba, GTX 1650 chỉ 4GB VRAM bị OOM — em dùng mixed precision fp16 và batch 12. "
        "Bốn, EasyOCR cold-start 20 giây — em pre-warm trong FastAPI lifespan(). "
        "Năm, Gemini quota free 1500 req/day hết khi demo — em build dual backend "
        "với SDK direct và OpenAI-compat proxy fallback. "
        "Sáu, class imbalance 200k chỉ 457 ảnh trong khi class khác 1800 — "
        "em dùng WeightedRandomSampler và class_weight 1.3×. "
        "Tổng cộng em chạy 7 sprints, khoảng 200 commits, refactor 4 god-classes về dưới 300 dòng."
    )


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 19 — SWOT ANALYSIS (NEW)
# ═══════════════════════════════════════════════════════════════════════════

def slide_swot_analysis(prs):
    """SWOT + AI reality + Unity flexibility — Phân tích thực tế."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(s, "#FFFFFF", "#EFF6FF", angle=315)
    add_section_label(s, 17, "Phân tích thực tế")
    add_logo(s)
    add_title(s, "Hệ thống có thực sự sẵn sàng cho thực tế?")

    # 2x2 SWOT grid
    cw = Inches(6.05)
    ch = Inches(2.0)
    gap = Inches(0.15)
    g_left = Inches(0.5)
    g_top = Inches(2.3)
    swot = [
        ("ĐIỂM MẠNH", GREEN, [
            "Tự động 100% — không cần bảo vệ thường trực",
            "AI accuracy cao: plate 94.5%, banknote 98.22%, chatbot >93%",
            "Open-source full stack — không lock-in vendor",
            "Pluggable IoT + Unity sim — test E2E mà chưa cần hardware",
        ]),
        ("ĐIỂM YẾU", RED, [
            "Chưa pilot trên bãi thực tế với user thật",
            "AI edge cases: tiền cũ/rách, biển bẩn, đêm tối quá",
            "Cần internet ổn định — mất net = offline mode chưa hoàn thiện",
            "Hardware ESP32 chưa deploy thật (chỉ simulator Unity)",
        ]),
        ("CƠ HỘI", BLUE, [
            "Smart parking VN còn sơ khai — open market",
            "Văn hóa cash payment VN phù hợp banknote AI",
            "Cloud + 5G giảm cost edge compute theo thời gian",
            "Multi-tenant marketplace — nhân rộng bãi pilot",
        ]),
        ("THÁCH THỨC", ORANGE, [
            "Pháp lý: biển số xe là personal data — cần consent",
            "Cạnh tranh từ Hikvision, Dahua, ParkingDesk",
            "User adoption — bãi truyền thống khó đổi sang tự động",
            "Bão táp/ăn cắp camera — hardware reliability",
        ]),
    ]
    for i, (title, c, items) in enumerate(swot):
        col = i % 2
        row = i // 2
        left = g_left + (cw + gap) * col
        top = g_top + (ch + gap) * row
        add_rect(s, left, top, cw, ch, WHITE, line=GRAY_2)
        add_rect(s, left, top, cw, Inches(0.4), c)
        add_text(s, left + Inches(0.2), top + Inches(0.06),
                 cw - Inches(0.3), Inches(0.3),
                 title, size=14, bold=True, color=WHITE, font=FONT_BOLD)
        body_text = "\n".join("• " + x for x in items)
        add_text(s, left + Inches(0.25), top + Inches(0.5),
                 cw - Inches(0.45), ch - Inches(0.55),
                 body_text, size=11, color=SLATE, line_spacing=1.35)

    # Bottom — 3 key questions answered
    add_rect(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5), BLUE_DEEP)
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
             "AI đủ tốt cho production? — Có với manual override fallback  ·  Unity flexibility? — Config JSON regen scene",
             size=12, bold=True, color=ORANGE_SOFT,
             align=PP_ALIGN.CENTER, font=FONT_BOLD,
             anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, 18)
    add_notes(s,
        "Em phân tích thực tế hệ thống qua mô hình SWOT. "
        "Điểm mạnh: tự động hóa 100% không cần bảo vệ thường trực, AI accuracy "
        "cao trên test set, open-source không lock-in vendor, có Unity simulator để "
        "test E2E mà chưa cần deploy hardware thật. "
        "Điểm yếu: chưa pilot trên bãi thực với user thật, AI gặp edge case khi tiền "
        "cũ rách hay biển bẩn, dependence vào internet ổn định, hardware ESP32 chưa "
        "deploy thật mà chỉ simulator Unity. "
        "Cơ hội: smart parking Việt Nam còn rất sơ khai, văn hóa cash payment phù hợp "
        "với banknote AI của em, cloud computing giảm cost theo thời gian. "
        "Thách thức: pháp lý về biển số là personal data cần consent, cạnh tranh từ "
        "vendor lớn như Hikvision Dahua, và user adoption khó vì bãi truyền thống "
        "không dễ chuyển đổi. "
        "Về câu hỏi quan trọng: AI có đủ tốt cho production chưa? Em trả lời CÓ — "
        "với điều kiện có manual override cho edge case. Banknote 99.5% precision-at-accept "
        "đã đủ ngưỡng bank-grade. Plate 94.5% e2e + Levenshtein tolerance đảm bảo false "
        "accept < 0.5%. Unity flexibility — chỉ cần thay JSON config, ParkingLotGenerator "
        "tự regen scene, đã hỗ trợ 4 mô hình bãi khác nhau."
    )




# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 20 — ROADMAP + BUSINESS CASE
# ═══════════════════════════════════════════════════════════════════════════

def slide_19(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(s, "#FFFFFF", "#EFF6FF", angle=315)
    add_section_label(s, 18, "Phát triển & Hiệu quả")
    add_logo(s)
    add_title(s, "Hướng phát triển  +  Business Case")

    # 6 roadmap cards (3x2)
    items = [
        ("01", "Hoạt động ngoại tuyến", "Local fallback queue khi mất internet", BLUE),
        ("02", "Cảm biến nâng cao", "LiDAR cho slot detection chính xác 100%", GREEN),
        ("03", "Sạc xe điện EV", "Detect xe điện + tính tiền sạc tự động", ORANGE),
        ("04", "Dự đoán thông minh", "ML dự đoán giờ bãi đầy → recommend đặt", PURPLE),
        ("05", "Mobile native app", "iOS/Android · push notification realtime", RED),
        ("06", "Multi-tenant", "Mỗi chủ bãi 1 tenant · platform Grab Park", SLATE_2),
    ]
    cw = Inches(4.0)
    ch = Inches(1.35)
    gap = Inches(0.13)
    grid_w = cw * 3 + gap * 2
    g_left = (SLIDE_W - grid_w) / 2
    g_top = Inches(2.4)
    for i, (num, title, desc, c) in enumerate(items):
        col = i % 3
        row = i // 3
        left = g_left + (cw + gap) * col
        top = g_top + (ch + gap) * row
        add_rect(s, left, top, cw, ch, WHITE, line=GRAY_2)
        add_rect(s, left, top, Inches(0.12), ch, c)
        # Number circle
        add_rect(s, left + cw - Inches(0.55), top + Inches(0.1),
                 Inches(0.4), Inches(0.4), c, shape=MSO_SHAPE.OVAL)
        add_text(s, left + cw - Inches(0.55), top + Inches(0.1),
                 Inches(0.4), Inches(0.4), num,
                 size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_BOLD,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, left + Inches(0.25), top + Inches(0.18),
                 cw - Inches(0.85), Inches(0.45),
                 title, size=15, bold=True, color=c, font=FONT_BOLD)
        add_text(s, left + Inches(0.25), top + Inches(0.7),
                 cw - Inches(0.4), ch - Inches(0.8),
                 desc, size=12, color=SLATE_2, line_spacing=1.35)

    # Business case panel — above footer
    bc_top = Inches(5.85)
    add_rect(s, Inches(0.5), bc_top, Inches(12.3), Inches(1.05), BLUE_DEEP)
    add_text(s, Inches(0.7), bc_top + Inches(0.08), Inches(12), Inches(0.3),
             "BUSINESS CASE — bãi 50 chỗ trong 1 năm",
             size=11, bold=True, color=ORANGE, font=FONT_BOLD)
    bc = [("~26.5tr", "Phần cứng (1 lần)"),
          ("8tr", "Vận hành/tháng"),
          ("35tr", "Doanh thu/tháng"),
          ("~1", "Tháng payback")]
    bw = Inches(2.9)
    bg_total = bw * 4 + Inches(0.1) * 3
    bcl = (SLIDE_W - bg_total) / 2
    for i, (v, lb) in enumerate(bc):
        l = bcl + (bw + Inches(0.1)) * i
        add_text(s, l, bc_top + Inches(0.4), bw, Inches(0.4),
                 v, size=22, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_BOLD)
        add_text(s, l, bc_top + Inches(0.78), bw, Inches(0.25),
                 lb, size=11, color=BLUE_SOFT, italic=True,
                 align=PP_ALIGN.CENTER)
    add_footer(s, 19)
    add_notes(s,
        "Em có 6 hướng phát triển sau khi tốt nghiệp. "
        "Một, Hoạt động ngoại tuyến — local fallback queue khi mất internet, sync khi có lại. "
        "Hai, Cảm biến nâng cao — tích hợp LiDAR cho slot detection chính xác 100%. "
        "Ba, Sạc xe điện — detect biển xanh và tính tiền sạc tự động. "
        "Bốn, Dự đoán thông minh — mô hình ML LSTM hoặc Prophet dự đoán giờ bãi đầy. "
        "Năm, Mobile native app — React Native với push notification realtime. "
        "Sáu, Multi-tenant marketplace — mô hình Grab Park, mỗi chủ bãi 1 tenant. "
        "Về business case cho bãi 50 chỗ: phần cứng ~26.5 triệu một lần (giá thực tế), vận hành 8 triệu/tháng (không cần bảo vệ thường trực), "
        "doanh thu ước tính 35 triệu/tháng (50 chỗ × 50% occupancy × 50k/ngày), payback period chỉ ~1 tháng. "
        "Đây là tính cho bãi prime nội đô với 70% occupancy. Bãi ngoại ô có thể payback 4-6 tháng."
    )


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 20 — ĐÓNG GÓP & KHÁC BIỆT
# ═══════════════════════════════════════════════════════════════════════════

def slide_20(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(s, "#FFFFFF", "#EFF6FF", angle=315)
    add_section_label(s, 19, "Đóng góp & Khác biệt")
    add_logo(s)
    add_title(s, "5 Unique Selling Points")

    # 5 USPs
    usps = [
        ("AI nhận diện 9 mệnh giá tiền VN", "Bank-grade pipeline 3-stage · DUY NHẤT TẠI VN", ORANGE),
        ("Unity Digital Twin", "Procedural generation · 8 cameras · 4 ESP32 simulator", GREEN),
        ("Microservices Open-Source Stack", "Đa ngôn ngữ · không lock-in vendor", BLUE),
        ("Chatbot LLM + RAG", "16 intents · Hybrid Confidence · citation từ docs", PURPLE),
        ("Precision-First AI", "Banknote ≥99.5% precision · Plate Levenshtein ≤ 1", RED),
    ]
    u_top = Inches(2.35)
    u_h = Inches(0.78)
    u_gap = Inches(0.06)
    for i, (title, desc, c) in enumerate(usps):
        top = u_top + (u_h + u_gap) * i
        add_rect(s, Inches(0.5), top, Inches(12.3), u_h, GRAY_SOFT)
        # Number circle
        add_rect(s, Inches(0.7), top + Inches(0.13),
                 Inches(0.6), Inches(0.6), c, shape=MSO_SHAPE.OVAL)
        add_text(s, Inches(0.7), top + Inches(0.13),
                 Inches(0.6), Inches(0.6), str(i+1),
                 size=22, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_BOLD,
                 anchor=MSO_ANCHOR.MIDDLE)
        # Title
        add_text(s, Inches(1.55), top + Inches(0.1),
                 Inches(11), Inches(0.4),
                 title, size=18, bold=True, color=c, font=FONT_BOLD)
        # Desc
        add_text(s, Inches(1.55), top + Inches(0.5),
                 Inches(11), Inches(0.35),
                 desc, size=13, color=SLATE_2, italic=True)

    # Bottom quote — above footer
    add_rect(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4), BLUE)
    add_text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4),
             '"Không phải làm tốt 1 thứ, mà làm tốt 1 hệ thống tích hợp"',
             size=14, italic=True, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font=FONT_BOLD,
             anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, 20)
    add_notes(s,
        "Đóng góp khoa học và điểm khác biệt của ParkSmart gồm 5 unique selling points. "
        "Một, AI nhận diện 9 mệnh giá tiền Việt Nam — duy nhất tại Việt Nam, "
        "với bank-grade pipeline 3-stage gồm EfficientNetV2-S Classifier, Siamese Network "
        "và OneClass SVM. Chưa bãi xe nào hỗ trợ cash payment qua AI. "
        "Hai, Unity Digital Twin — procedural lot generation, 8 virtual cameras, "
        "4 ESP32 simulator — không có project KLTN VN nào tương tự. "
        "Ba, Microservices Open-Source Stack — đa ngôn ngữ Python, Go, JavaScript, C#, "
        "không lock-in vendor, ai cũng có thể fork và deploy. "
        "Bốn, Chatbot LLM với RAG — không phải FAQ template bot, mà thực sự hiểu ngữ cảnh, "
        "có 16 intents, Hybrid Confidence và citation từ knowledge base thật. "
        "Năm, Precision-First AI cho giao dịch tài chính — banknote precision-at-accept "
        "≥ 99.5%, plate Levenshtein ≤ 1, false accept rate < 0.5%. "
        "Triết lý của em là: không phải làm tốt 1 thứ, mà làm tốt 1 hệ thống tích hợp."
    )


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 21 — CẢM ƠN / Q&A
# ═══════════════════════════════════════════════════════════════════════════

def slide_21(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(s, "#0A1330", "#1E3A5F", angle=315)
    add_rect(s, Inches(0), Inches(3.0), Inches(0.5), Inches(2.0), ORANGE)
    add_rect(s, Inches(12.83), Inches(2.5), Inches(0.5), Inches(2.0), ORANGE)
    # University logo top-left for KLTN consistency
    add_logo(s, mode="large")
    add_text(s, Inches(1.6), Inches(0.45), Inches(7), Inches(0.4),
             "TRƯỜNG ĐẠI HỌC MỞ THÀNH PHỐ HỒ CHÍ MINH",
             size=12, bold=True, color=ORANGE)
    add_text(s, Inches(1.6), Inches(0.8), Inches(7), Inches(0.35),
             "KHOA CÔNG NGHỆ THÔNG TIN",
             size=10, color=ORANGE_SOFT)

    add_text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(1.2),
             "CẢM ƠN HỘI ĐỒNG",
             size=64, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font=FONT_BOLD,
             line_spacing=1.0)
    add_text(s, Inches(0.7), Inches(3.3), Inches(12), Inches(0.7),
             "Q & A  ·  Mời hội đồng đặt câu hỏi",
             size=24, color=ORANGE_SOFT, italic=True,
             align=PP_ALIGN.CENTER)
    add_rect(s, (SLIDE_W - Inches(2.5)) / 2, Inches(4.3),
             Inches(2.5), Inches(0.05), ORANGE)

    add_text(s, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.4),
             "DEMO LIVE", size=12, bold=True, color=ORANGE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.5),
             "https://parksmart.ghepdoicaulong.shop",
             size=22, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font=FONT_BOLD)

    add_text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.4),
             "LIÊN HỆ", size=12, bold=True, color=ORANGE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.4),
             "Nguyễn Hải Minh  ·  MSSV 2251012093",
             size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font=FONT_BOLD)
    add_text(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.4),
             "dang.nguyenhai2k2@gmail.com",
             size=14, color=BLUE_SOFT,
             align=PP_ALIGN.CENTER)

    add_notes(s,
        "Em xin kết thúc phần trình bày tại đây. Cảm ơn hội đồng đã lắng nghe. "
        "Em sẵn sàng nhận câu hỏi từ hội đồng. Em cũng có demo live tại "
        "parksmart chấm ghepdoicaulong chấm shop nếu hội đồng muốn xem trực tiếp."
    )


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 22 — Q&A APPENDIX (Technical)
# ═══════════════════════════════════════════════════════════════════════════

def slide_22(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(s, "#FFFFFF", "#EFF6FF", angle=315)
    add_section_label(s, 20, "Phụ lục — Câu hỏi (1/2)", accent=PURPLE)
    add_logo(s)
    add_title(s, "Câu hỏi về AI và Kiến trúc")

    qs = [
        "Cấu trúc tổng thể hệ thống ParkSmart như thế nào?",
        "Tại sao chia microservice thay vì monolith?",
        "Mô hình nhận diện tiền hoạt động ra sao?",
        "Confusion matrix banknote — class nào nhầm?",
        "Làm sao nâng cao thuật toán banknote thêm?",
        "AI plate xử lý biển nghiêng/mờ/đêm thế nào?",
        "RAG hoạt động ra sao trong chatbot?",
        "Tại sao chọn ESP32 không phải Pi/Arduino?",
        "Đã deploy hardware ESP32 thật chưa?",
        "Multi-floor NavMesh Unity hoạt động chưa?",
    ]
    q_top = Inches(2.4)
    q_h = Inches(0.42)
    cw = Inches(6.0)
    g = Inches(0.3)
    for i, q in enumerate(qs):
        col = i // 5
        row = i % 5
        left = Inches(0.5) + (cw + g) * col
        top = q_top + (q_h + Inches(0.08)) * row
        add_rect(s, left, top, cw, q_h, GRAY_SOFT)
        # Number badge
        add_rect(s, left + Inches(0.1), top + Inches(0.05),
                 Inches(0.4), q_h - Inches(0.1), PURPLE, shape=MSO_SHAPE.OVAL)
        add_text(s, left + Inches(0.1), top + Inches(0.05),
                 Inches(0.4), q_h - Inches(0.1),
                 str(i+1), size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_BOLD,
                 anchor=MSO_ANCHOR.MIDDLE)
        # Question
        add_text(s, left + Inches(0.6), top,
                 cw - Inches(0.7), q_h,
                 q, size=13, color=SLATE,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)

    add_text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4),
             "Slide backup — chỉ mở khi panel hỏi đúng câu hỏi tương ứng",
             size=12, italic=True, color=GRAY,
             align=PP_ALIGN.CENTER)
    add_footer(s, 22)
    add_notes(s,
        "Đây là slide phụ lục — em chỉ mở khi hội đồng hỏi đúng câu hỏi tương ứng. "
        "10 câu hỏi đầu tiên về AI và kiến trúc. Em đã chuẩn bị câu trả lời cho từng câu, "
        "có thể giải thích chi tiết nếu hội đồng yêu cầu."
    )


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 23 — Q&A APPENDIX (Practical)
# ═══════════════════════════════════════════════════════════════════════════

def slide_23(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(s, "#FFFFFF", "#EFF6FF", angle=315)
    add_section_label(s, 21, "Phụ lục — Câu hỏi (2/2)", accent=PURPLE)
    add_logo(s)
    add_title(s, "Câu hỏi về Thực tế và Khác biệt")

    qs = [
        "Xe đến trễ/sớm so với booking thì sao?",
        "Booking confirmed nhưng user không đến?",
        "Bãi đầy 100% xử lý sao?",
        "Mất internet / ESP32 mất kết nối?",
        "Đã khảo sát bãi giữ xe thực tế ở đâu?",
        "Cơ chế chứng minh quy trình hoạt động đúng?",
        "Có học hỏi từ smart parking Trung Quốc không?",
        "Cách dựng Unity tốt nhất + cách trình bày?",
        "Chi phí triển khai thực tế chi tiết?",
        "ParkSmart khác biệt gì với smart parking khác?",
    ]
    q_top = Inches(2.4)
    q_h = Inches(0.42)
    cw = Inches(6.0)
    g = Inches(0.3)
    for i, q in enumerate(qs):
        col = i // 5
        row = i % 5
        left = Inches(0.5) + (cw + g) * col
        top = q_top + (q_h + Inches(0.08)) * row
        add_rect(s, left, top, cw, q_h, GRAY_SOFT)
        add_rect(s, left + Inches(0.1), top + Inches(0.05),
                 Inches(0.4), q_h - Inches(0.1), PURPLE, shape=MSO_SHAPE.OVAL)
        add_text(s, left + Inches(0.1), top + Inches(0.05),
                 Inches(0.4), q_h - Inches(0.1),
                 str(i+11), size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_BOLD,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, left + Inches(0.6), top,
                 cw - Inches(0.7), q_h,
                 q, size=13, color=SLATE,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)

    add_text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4),
             "Slide backup — chỉ mở khi panel hỏi đúng câu hỏi tương ứng",
             size=12, italic=True, color=GRAY,
             align=PP_ALIGN.CENTER)
    add_footer(s, 23)
    add_notes(s,
        "10 câu hỏi tiếp theo về thực tế áp dụng và điểm khác biệt. "
        "Em đã chuẩn bị câu trả lời cho từng câu — về cách xử lý xe đến trễ/sớm, "
        "bãi đầy, mất internet, khảo sát bãi thực tế, học hỏi smart parking Trung Quốc, "
        "chi phí triển khai và điểm khác biệt với competitor."
    )


# ═══════════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [slide_01, slide_02, slide_03, slide_04, slide_05,
                slide_06, slide_07, slide_08, slide_09, slide_10,
                slide_11, slide_12, slide_13, slide_14, slide_15,
                slide_16, slide_18, slide_swot_analysis, slide_19, slide_20,
                slide_21, slide_22, slide_23]
    for b in builders:
        b(prs)

    out = Path(__file__).parent / "ParkSmart_KLTN_Slides.pptx"
    try:
        prs.save(out)
    except PermissionError:
        from datetime import datetime
        suffix = datetime.now().strftime("%H%M%S")
        out = Path(__file__).parent / f"ParkSmart_KLTN_Slides_v8_{suffix}.pptx"
        prs.save(out)
        print("[!] File gốc đang mở — đã lưu suffix mới")
    print(f"OK saved: {out.name}  ({out.stat().st_size // 1024} KB)  ·  {len(builders)} slides")
    return out


if __name__ == "__main__":
    build()
